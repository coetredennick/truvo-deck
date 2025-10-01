from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import glob

prs = Presentation()
# 16:9 canvas
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

for img_path in sorted(glob.glob("slide-*.png")):
    print(f"Adding {img_path}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    img = Image.open(img_path)
    iw, ih = img.size
    # scale to cover full slide
    slide.shapes.add_picture(
        img_path, 0, 0, width=prs.slide_width, height=prs.slide_height
    )

prs.save("truvo-pitch-deck.pptx")
print("Done! PowerPoint saved as truvo-pitch-deck.pptx")
